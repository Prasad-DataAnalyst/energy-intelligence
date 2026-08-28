"""
data/ingest_documents.py — Document Ingestion Engine
=====================================================
Watches INPUT_DIRS for new files, extracts text from every supported format,
calls Claude Haiku (or heuristic fallback) to extract video-ready insights,
and persists results to Memory for the content pipeline to consume.
"""
import json
import logging
import os
import re
import threading
import time
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional

log = logging.getLogger("brain.ingest_documents")

HERE = Path(__file__).parent

INPUT_DIRS = {
    "blogs":       HERE / "input" / "blogs",
    "articles":    HERE / "input" / "articles",
    "research":    HERE / "input" / "research",
    "reports":     HERE / "input" / "reports",
    "books":       HERE / "input" / "books",
    "scripts":     HERE / "input" / "scripts",
    "notes":       HERE / "input" / "notes",
    "data":        HERE / "input" / "data",
    "archives":    HERE / "input" / "archives",
    "private":     HERE / "input" / "private",
    "newsletters": HERE / "input" / "newsletters",
}

MAX_TEXT_CHARS = 15_000
CLAUDE_MODEL   = "claude-haiku-4-5-20251001"


# ─────────────────────────────────────────────────────────────────────────────
# Class 1: DocumentExtractor
# ─────────────────────────────────────────────────────────────────────────────

class DocumentExtractor:
    """Extract raw text from any supported document format."""

    def extract(self, file_path: Path) -> str:
        """Dispatch to the correct extractor based on file suffix."""
        suffix = file_path.suffix.lower()
        dispatch = {
            ".pdf":  self._extract_pdf,
            ".docx": self._extract_docx,
            ".xlsx": self._extract_xlsx,
            ".xls":  self._extract_xlsx,
            ".csv":  self._extract_csv,
            ".epub": self._extract_epub,
            ".html": self._extract_html,
            ".htm":  self._extract_html,
            ".txt":  self._extract_txt,
            ".md":   self._extract_txt,
            ".pptx": self._extract_pptx,
            ".json": self._extract_json,
        }
        handler = dispatch.get(suffix, self._extract_txt)
        try:
            text = handler(file_path)
        except Exception as exc:
            log.warning("extract(%s) unhandled error: %s", file_path.name, exc)
            text = ""
        return (text or "")[:MAX_TEXT_CHARS]

    def _extract_pdf(self, path: Path) -> str:
        # Try PyMuPDF first
        try:
            import fitz  # type: ignore
            doc  = fitz.open(str(path))
            text = "\n".join(page.get_text() for page in doc)
            doc.close()
            if text.strip():
                return text
        except ImportError:
            log.debug("fitz not available, trying pdfplumber")
        except Exception as exc:
            log.debug("fitz failed for %s: %s", path.name, exc)

        # Fallback to pdfplumber
        try:
            import pdfplumber  # type: ignore
            with pdfplumber.open(str(path)) as pdf:
                text = "\n".join(
                    (page.extract_text() or "") for page in pdf.pages
                )
            if text.strip():
                return text
        except ImportError:
            log.debug("pdfplumber not available, falling back to raw read")
        except Exception as exc:
            log.debug("pdfplumber failed for %s: %s", path.name, exc)

        # Last resort: raw bytes, strip non-printable
        try:
            raw = path.read_bytes()
            text = raw.decode("utf-8", errors="replace")
            # Keep only printable ASCII-ish content
            text = re.sub(r"[^\x20-\x7e\n\r\t]", " ", text)
            return re.sub(r" {4,}", "  ", text)
        except Exception as exc:
            log.debug("raw PDF read failed for %s: %s", path.name, exc)
            return ""

    def _extract_docx(self, path: Path) -> str:
        try:
            import docx  # type: ignore
            doc   = docx.Document(str(path))
            paras = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paras)
        except ImportError:
            log.debug("python-docx not available")
            return ""
        except Exception as exc:
            log.debug("docx extraction failed for %s: %s", path.name, exc)
            return ""

    def _extract_xlsx(self, path: Path) -> str:
        lines: List[str] = []

        # openpyxl primary
        try:
            import openpyxl  # type: ignore
            wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
            for sheet in wb.worksheets:
                lines.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in cells):
                        lines.append("\t".join(cells))
            wb.close()
        except ImportError:
            log.debug("openpyxl not available")
        except Exception as exc:
            log.debug("openpyxl failed for %s: %s", path.name, exc)

        # pandas describe as supplement
        try:
            import pandas as pd  # type: ignore
            df_dict = pd.read_excel(str(path), sheet_name=None)
            for name, df in df_dict.items():
                desc = df.describe(include="all").to_string()
                lines.append(f"\n[Stats for sheet {name}]\n{desc}")
        except Exception:
            pass

        return "\n".join(lines)

    def _extract_csv(self, path: Path) -> str:
        try:
            import pandas as pd  # type: ignore
            df    = pd.read_csv(str(path), on_bad_lines="skip", nrows=500)
            head  = df.head(20).to_string(index=False)
            desc  = df.describe(include="all").to_string()
            cols  = f"Columns ({len(df.columns)}): {', '.join(str(c) for c in df.columns)}"
            rows  = f"Total rows: {len(df)}"
            return f"{cols}\n{rows}\n\n{head}\n\n[Statistics]\n{desc}"
        except ImportError:
            log.debug("pandas not available, falling back to text read")
            return self._extract_txt(path)
        except Exception as exc:
            log.debug("csv extraction failed for %s: %s", path.name, exc)
            return self._extract_txt(path)

    def _extract_epub(self, path: Path) -> str:
        try:
            import ebooklib  # type: ignore
            from ebooklib import epub as _epub
            try:
                from bs4 import BeautifulSoup as _BS  # type: ignore
                _use_bs = True
            except ImportError:
                _use_bs = False

            book  = _epub.read_epub(str(path))
            parts: List[str] = []
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    raw = item.get_content().decode("utf-8", errors="replace")
                    if _use_bs:
                        soup  = _BS(raw, "html.parser")
                        parts.append(soup.get_text(separator="\n"))
                    else:
                        parts.append(re.sub(r"<[^>]+>", " ", raw))
            return "\n\n".join(parts)
        except ImportError:
            log.debug("ebooklib not available")
            return ""
        except Exception as exc:
            log.debug("epub extraction failed for %s: %s", path.name, exc)
            return ""

    def _extract_html(self, path: Path) -> str:
        try:
            from bs4 import BeautifulSoup  # type: ignore
            html = path.read_text(encoding="utf-8", errors="replace")
            soup = BeautifulSoup(html, "html.parser")
            tags = soup.find_all(["p", "h1", "h2", "h3", "h4", "h5", "h6", "li"])
            return "\n".join(t.get_text(separator=" ", strip=True) for t in tags)
        except ImportError:
            log.debug("bs4 not available, stripping HTML manually")
            raw = path.read_text(encoding="utf-8", errors="replace")
            return re.sub(r"<[^>]+>", " ", raw)
        except Exception as exc:
            log.debug("html extraction failed for %s: %s", path.name, exc)
            return ""

    def _extract_txt(self, path: Path) -> str:
        try:
            return path.read_text(encoding="utf-8", errors="replace")
        except Exception as exc:
            log.debug("txt extraction failed for %s: %s", path.name, exc)
            return ""

    def _extract_pptx(self, path: Path) -> str:
        try:
            from pptx import Presentation  # type: ignore
            prs   = Presentation(str(path))
            parts: List[str] = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text: List[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_text.append(t)
                if slide_text:
                    parts.append(f"[Slide {i}]\n" + "\n".join(slide_text))
            return "\n\n".join(parts)
        except ImportError:
            log.debug("python-pptx not available")
            return ""
        except Exception as exc:
            log.debug("pptx extraction failed for %s: %s", path.name, exc)
            return ""

    def _extract_json(self, path: Path) -> str:
        try:
            raw  = path.read_text(encoding="utf-8", errors="replace")
            data = json.loads(raw)
            return json.dumps(data, indent=2)
        except Exception as exc:
            log.debug("json extraction failed for %s: %s", path.name, exc)
            return self._extract_txt(path)


# ─────────────────────────────────────────────────────────────────────────────
# Class 2: InsightExtractor
# ─────────────────────────────────────────────────────────────────────────────

class InsightExtractor:
    """Use Claude Haiku (or heuristics) to extract video-ready insights."""

    _SURPRISE_WORDS = {
        "surprising", "shocking", "unexpected", "revealed", "discovered",
        "first time", "never before", "unprecedented", "extraordinary",
        "remarkable", "astonishing", "staggering", "jaw-dropping",
    }

    def __init__(self):
        self._claude_key = os.getenv("ANTHROPIC_API_KEY", "")

    def extract_insights(
        self, text: str, source_type: str, filename: str
    ) -> List[Dict]:
        """Extract insights via Claude Haiku; fall back to heuristics."""
        if not self._claude_key:
            log.debug("No ANTHROPIC_API_KEY — using heuristic insights")
            return self._heuristic_insights(text, source_type)

        prompt = (
            f"Extract 5-8 key insights from this {source_type} document '{filename}'. "
            "Return JSON array: "
            '[{"text": str, "category": str, "emotion": str, '
            '"surprise_score": 1-100, "video_potential": 1-100, '
            '"key_stat": str or null, "angle": str}]. '
            "Focus on: surprising statistics, counterintuitive findings, "
            "controversial claims, human interest angles."
        )
        snippet = text[:12000]
        try:
            import anthropic  # type: ignore
            client = anthropic.Anthropic(api_key=self._claude_key)
            msg = client.messages.create(
                model=CLAUDE_MODEL,
                max_tokens=1500,
                messages=[
                    {
                        "role": "user",
                        "content": f"{prompt}\n\nDOCUMENT:\n{snippet}",
                    }
                ],
            )
            raw = msg.content[0].text.strip()
            # Strip markdown code fences if present
            raw = re.sub(r"^```(?:json)?\s*", "", raw, flags=re.MULTILINE)
            raw = re.sub(r"\s*```$", "", raw, flags=re.MULTILINE)
            raw = raw.strip()
            insights = json.loads(raw)
            if not isinstance(insights, list):
                raise ValueError("Expected JSON array")
            validated: List[Dict] = []
            for item in insights:
                if not isinstance(item, dict):
                    continue
                validated.append({
                    "text":            str(item.get("text", "")),
                    "category":        str(item.get("category", source_type)),
                    "emotion":         str(item.get("emotion", "neutral")),
                    "surprise_score":  int(item.get("surprise_score", 50)),
                    "video_potential": int(item.get("video_potential", 50)),
                    "key_stat":        item.get("key_stat"),
                    "angle":           str(item.get("angle", "")),
                    "source_type":     source_type,
                    "filename":        filename,
                    "extracted_at":    datetime.now(timezone.utc).isoformat(),
                })
            return validated
        except json.JSONDecodeError as exc:
            log.warning("Claude returned non-JSON for %s: %s", filename, exc)
            return self._heuristic_insights(text, source_type)
        except Exception as exc:
            log.warning("Claude insight extraction failed for %s: %s", filename, exc)
            return self._heuristic_insights(text, source_type)

    def _heuristic_insights(self, text: str, source_type: str) -> List[Dict]:
        """Extract insights without API using text pattern matching."""
        now      = datetime.now(timezone.utc).isoformat()
        insights: List[Dict] = []

        # Tokenise into sentences
        sentences = re.split(r"(?<=[.!?])\s+", text)

        for sent in sentences:
            sent = sent.strip()
            if len(sent) < 30 or len(sent) > 400:
                continue

            # Stats insights: sentences containing numbers / percentages
            if re.search(r"\b\d[\d,.]*\s*%|\b\d{2,}", sent):
                stat_match = re.search(r"(\d[\d,.]*\s*%?)", sent)
                insights.append({
                    "text":            sent,
                    "category":        "statistic",
                    "emotion":         "informative",
                    "surprise_score":  60,
                    "video_potential": 70,
                    "key_stat":        stat_match.group(1) if stat_match else None,
                    "angle":           "data-driven revelation",
                    "source_type":     source_type,
                    "filename":        "",
                    "extracted_at":    now,
                })
                if len(insights) >= 10:
                    break
                continue

            # Surprise insights
            sent_lower = sent.lower()
            if any(word in sent_lower for word in self._SURPRISE_WORDS):
                insights.append({
                    "text":            sent,
                    "category":        "revelation",
                    "emotion":         "surprise",
                    "surprise_score":  80,
                    "video_potential": 80,
                    "key_stat":        None,
                    "angle":           "shocking discovery",
                    "source_type":     source_type,
                    "filename":        "",
                    "extracted_at":    now,
                })
                if len(insights) >= 10:
                    break
                continue

            # Question insights
            if sent.endswith("?"):
                insights.append({
                    "text":            sent,
                    "category":        "question",
                    "emotion":         "curiosity",
                    "surprise_score":  50,
                    "video_potential": 60,
                    "key_stat":        None,
                    "angle":           "audience question hook",
                    "source_type":     source_type,
                    "filename":        "",
                    "extracted_at":    now,
                })
                if len(insights) >= 10:
                    break

        return insights[:10]

    def score_video_potential(self, insight: Dict) -> float:
        """Weighted average of surprise_score (40%) and video_potential (60%)."""
        surprise = float(insight.get("surprise_score", 50))
        vp       = float(insight.get("video_potential", 50))
        return round(surprise * 0.4 + vp * 0.6, 2)


# ─────────────────────────────────────────────────────────────────────────────
# Class 3: FolderWatcher
# ─────────────────────────────────────────────────────────────────────────────

class FolderWatcher:
    """Scan INPUT_DIRS for new documents and trigger ingestion."""

    def __init__(self, memory=None):
        self._memory    = memory
        self._extractor = DocumentExtractor()
        self._insights  = InsightExtractor()
        self._stop      = threading.Event()

        # Ensure all input directories exist
        for key, path in INPUT_DIRS.items():
            try:
                path.mkdir(parents=True, exist_ok=True)
            except Exception as exc:
                log.warning("Could not create dir %s: %s", path, exc)

    def get_supported_extensions(self) -> set:
        return {
            ".pdf", ".docx", ".xlsx", ".csv", ".epub",
            ".html", ".txt", ".md", ".pptx", ".json",
        }

    def scan_all_folders(self) -> List[Dict]:
        """Walk all INPUT_DIRS, process files modified in the last 24 hours."""
        cutoff    = datetime.now(timezone.utc) - timedelta(hours=24)
        processed: List[Dict] = []
        exts      = self.get_supported_extensions()

        for folder_key, folder_path in INPUT_DIRS.items():
            if not folder_path.exists():
                continue
            for file_path in folder_path.rglob("*"):
                if not file_path.is_file():
                    continue
                if file_path.suffix.lower() not in exts:
                    continue
                try:
                    mtime = datetime.fromtimestamp(
                        file_path.stat().st_mtime, tz=timezone.utc
                    )
                except OSError:
                    continue
                if mtime < cutoff:
                    continue
                if self._is_already_processed(file_path):
                    log.debug("Skipping already-processed: %s", file_path.name)
                    continue
                result = self.process_file(file_path, folder_key)
                if result:
                    processed.append(result)

        log.info("scan_all_folders: %d new files processed", len(processed))
        return processed

    def process_file(self, file_path: Path, folder_key: str) -> Dict:
        """Extract text + insights from one file, save to memory."""
        log.info("Processing %s (%s)", file_path.name, folder_key)

        # Extract text
        text = ""
        try:
            text = self._extractor.extract(file_path)
        except Exception as exc:
            log.warning("Text extraction failed for %s: %s", file_path.name, exc)

        if not text.strip():
            log.debug("No text extracted from %s", file_path.name)

        # Extract insights
        insights: List[Dict] = []
        try:
            insights = self._insights.extract_insights(
                text, folder_key, file_path.name
            )
        except Exception as exc:
            log.warning("Insight extraction failed for %s: %s", file_path.name, exc)

        # Score insights
        for ins in insights:
            ins["score"] = self._insights.score_video_potential(ins)
            ins["file_path"] = str(file_path)

        top_score = max((i["score"] for i in insights), default=0.0)

        # Persist to memory
        if self._memory is not None:
            try:
                self._memory.save_ingested_file({
                    "file_path":     str(file_path),
                    "filename":      file_path.name,
                    "folder":        folder_key,
                    "text_length":   len(text),
                    "insight_count": len(insights),
                    "top_score":     top_score,
                    "ingested_at":   datetime.now(timezone.utc).isoformat(),
                })
            except Exception as exc:
                log.warning("save_ingested_file failed for %s: %s", file_path.name, exc)

            for ins in insights:
                try:
                    self._memory.save_insight(ins)
                except Exception as exc:
                    log.warning("save_insight failed: %s", exc)

        return {
            "filename":      file_path.name,
            "folder":        folder_key,
            "insight_count": len(insights),
            "top_score":     top_score,
        }

    def _is_already_processed(self, file_path: Path) -> bool:
        """Return True if memory already has a record for this file path."""
        if self._memory is None:
            return False
        try:
            return self._memory.get_ingested_file(str(file_path)) is not None
        except Exception:
            return False

    def start_watching(self) -> threading.Thread:
        """Start background file-watching thread; use watchdog if available."""
        try:
            from watchdog.observers import Observer  # type: ignore
            from watchdog.events import FileSystemEventHandler  # type: ignore

            watcher = self

            class _Handler(FileSystemEventHandler):
                def on_created(self, event):
                    if event.is_directory:
                        return
                    path = Path(event.src_path)
                    if path.suffix.lower() not in watcher.get_supported_extensions():
                        return
                    if watcher._is_already_processed(path):
                        return
                    # Determine folder_key
                    folder_key = "unknown"
                    for key, folder_path in INPUT_DIRS.items():
                        try:
                            path.relative_to(folder_path)
                            folder_key = key
                            break
                        except ValueError:
                            pass
                    try:
                        watcher.process_file(path, folder_key)
                    except Exception as exc:
                        log.warning("watchdog handler error: %s", exc)

            observer = Observer()
            handler  = _Handler()
            for folder_path in INPUT_DIRS.values():
                if folder_path.exists():
                    observer.schedule(handler, str(folder_path), recursive=True)

            def _run():
                observer.start()
                log.info("Watchdog observer started")
                while not self._stop.is_set():
                    time.sleep(1)
                observer.stop()
                observer.join()
                log.info("Watchdog observer stopped")

            t = threading.Thread(target=_run, daemon=True, name="folder-watcher-watchdog")
            t.start()
            return t

        except ImportError:
            log.info("watchdog not available — falling back to polling every 300s")
            return self._start_polling()

    def _start_polling(self) -> threading.Thread:
        """Poll folders every 300 seconds when watchdog is unavailable."""
        def _poll():
            log.info("Polling watcher started (interval=300s)")
            while not self._stop.is_set():
                try:
                    self.scan_all_folders()
                except Exception as exc:
                    log.warning("Poll scan error: %s", exc)
                # Sleep in small increments so stop event is responsive
                for _ in range(300):
                    if self._stop.is_set():
                        break
                    time.sleep(1)
            log.info("Polling watcher stopped")

        t = threading.Thread(target=_poll, daemon=True, name="folder-watcher-poll")
        t.start()
        return t

    def stop_watching(self):
        """Signal background watcher to stop."""
        self._stop.set()


# ─────────────────────────────────────────────────────────────────────────────
# Class 4: DocumentIngestionEngine (orchestrator)
# ─────────────────────────────────────────────────────────────────────────────

class DocumentIngestionEngine:
    """Top-level orchestrator for the document ingestion pipeline."""

    def __init__(self, memory=None):
        if memory is None:
            try:
                from core.memory import get_memory  # type: ignore
                memory = get_memory()
            except Exception as exc:
                log.warning("Could not load memory: %s", exc)
        self._memory  = memory
        self._watcher = FolderWatcher(memory=self._memory)

    def run_full_scan(self) -> Dict:
        """Scan all folders, return summary statistics."""
        log.info("Starting full document scan…")
        results = self._watcher.scan_all_folders()

        total_insights = sum(r.get("insight_count", 0) for r in results)

        # Collect all insights to find top 5
        all_insights: List[Dict] = []
        if self._memory is not None:
            try:
                all_insights = self._memory.get_top_insights(n=5, min_score=0)
            except Exception as exc:
                log.debug("get_top_insights failed: %s", exc)

        # Fallback: derive top insights from scan results
        if not all_insights:
            flat: List[Dict] = []
            for r in results:
                flat.extend(r.get("insights", []))
            all_insights = sorted(flat, key=lambda x: x.get("score", 0), reverse=True)[:5]

        return {
            "files_processed": len(results),
            "total_insights":  total_insights,
            "top_insights":    all_insights[:5],
            "file_details":    results,
        }

    def ingest_single_file(self, file_path: str) -> Dict:
        """Process one specific file path and return summary."""
        path = Path(file_path)
        if not path.exists():
            log.error("File not found: %s", file_path)
            return {"error": f"File not found: {file_path}"}
        if not path.is_file():
            log.error("Path is not a file: %s", file_path)
            return {"error": f"Not a file: {file_path}"}

        # Determine folder key from path
        folder_key = "unknown"
        for key, folder_path in INPUT_DIRS.items():
            try:
                path.relative_to(folder_path)
                folder_key = key
                break
            except ValueError:
                pass

        return self._watcher.process_file(path, folder_key)

    def get_top_content_ideas(self, n: int = 10) -> List[Dict]:
        """Return top insights from memory ranked by video_potential score."""
        if self._memory is None:
            return []
        try:
            return self._memory.get_top_insights(n=n, min_score=0)
        except Exception as exc:
            log.warning("get_top_content_ideas failed: %s", exc)
            return []

    def print_report(self, result: Dict) -> None:
        """Print a formatted summary of an ingestion run."""
        sep = "=" * 60
        print(f"\n{sep}")
        print("  DOCUMENT INGESTION REPORT")
        print(sep)
        print(f"  Files processed : {result.get('files_processed', 0)}")
        print(f"  Total insights  : {result.get('total_insights', 0)}")
        if result.get("error"):
            print(f"  Error           : {result['error']}")
        top = result.get("top_insights", [])
        if top:
            print(f"\n  Top {len(top)} Insights (by score):")
            for i, ins in enumerate(top, 1):
                score = ins.get("score", ins.get("video_potential", 0))
                text  = (ins.get("text", "")[:80] + "…") if len(ins.get("text", "")) > 80 else ins.get("text", "")
                print(f"    {i}. [{score:.0f}] {text}")
        details = result.get("file_details", [])
        if details:
            print("\n  Processed Files:")
            for d in details:
                print(
                    f"    • {d.get('filename','?')} "
                    f"({d.get('folder','?')}) — "
                    f"{d.get('insight_count',0)} insights, "
                    f"top={d.get('top_score',0):.0f}"
                )
        print(f"{sep}\n")
